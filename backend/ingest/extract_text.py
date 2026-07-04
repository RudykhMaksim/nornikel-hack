"""Ограниченное извлечение текста по принципу best-effort из разнородного
реального корпуса.

Проектные ограничения (см. план): docx/docm/pptx/xlsx дёшевы (ZIP+XML,
UTF-8) и всегда полностью читаемы. PDF в этом корпусе имеют встроенные
шрифты без корректных ToUnicode CMaps для некоторых кодировок -- `pdftotext`
из xpdf молча теряет кириллический текст, поэтому для извлечения из PDF
вместо этого используется PyMuPDF (fitz), ограниченный числом страниц и
порогом размера файла (большие сканированные/вёрстанные выпуски
пропускаются для глубокого разбора и только каталогизируются как метаданные).
"""
import re
import zipfile
from pathlib import Path

MAX_PDF_PAGES = 20
MAX_CHARS = 20000
MAX_PDF_DEEP_SIZE_BYTES = 45 * 1024 * 1024  # PDF крупнее этого пропускаются при глубоком разборе

SUPPORTED_EXTS = {"docx", "docm", "pptx", "xlsx", "pdf"}


def extractable(ext: str, size_bytes: int) -> bool:
    ext = ext.lower()
    if ext not in SUPPORTED_EXTS:
        return False
    if ext == "pdf" and size_bytes > MAX_PDF_DEEP_SIZE_BYTES:
        return False
    return True


def _extract_docx(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text]
    for table in d.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    return "\n".join(parts)


def _extract_docx_fallback_zip(path: Path) -> str:
    """Извлечение из сырого XML -- используется, когда python-docx не может открыть .docm (с макросами)."""
    z = zipfile.ZipFile(path)
    xml = z.read("word/document.xml").decode("utf-8", errors="ignore")
    text = re.sub(r"<[^>]+>", " ", xml)
    return re.sub(r"\s+", " ", text).strip()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        parts.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text
            if note:
                parts.append(note)
    return "\n".join(parts)


def _extract_xlsx(path: Path, max_cells: int = 5000) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    count = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    parts.append(str(cell))
                    count += 1
                    if count >= max_cells:
                        wb.close()
                        return " ".join(parts)
    wb.close()
    return " ".join(parts)


def _extract_pdf(path: Path, max_pages: int = MAX_PDF_PAGES) -> str:
    import fitz
    doc = fitz.open(str(path))
    parts = []
    n = min(max_pages, doc.page_count)
    for i in range(n):
        parts.append(doc[i].get_text())
    doc.close()
    return "\n".join(parts)


def extract_text(path: Path, ext: str, max_chars: int = MAX_CHARS) -> str:
    ext = ext.lower()
    try:
        if ext == "docx":
            text = _extract_docx(path)
        elif ext == "docm":
            try:
                text = _extract_docx(path)
            except Exception:
                text = _extract_docx_fallback_zip(path)
        elif ext == "pptx":
            text = _extract_pptx(path)
        elif ext == "xlsx":
            text = _extract_xlsx(path)
        elif ext == "pdf":
            text = _extract_pdf(path)
        else:
            return ""
    except Exception:
        return ""
    return text[:max_chars]
